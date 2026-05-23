import os
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pdf(filepath):
    c = canvas.Canvas(filepath, pagesize=letter)
    c.setFont("Helvetica-Bold", 16)
    c.drawString(50, 750, "Reporte de Cambios: Problema de Roles de Usuario")
    
    c.setFont("Helvetica", 12)
    c.drawString(50, 710, "1. Descripción del Problema:")
    c.drawString(70, 690, "- Se creó un usuario (Lara Lala) que obtuvo el rol de Admin accidentalmente.")
    c.drawString(70, 670, "- El rol no se guardaba en la sesión al iniciar.")
    c.drawString(70, 650, "- La creación de usuarios no limitaba el rol de manera explícita en el controlador.")
    
    c.drawString(50, 610, "2. Solución Implementada:")
    c.drawString(70, 590, "- Se modificó usersController.js para forzar el rol 'user' por defecto.")
    c.drawString(70, 570, "- Se añadió lógica para permitir que solo administradores logueados asignen el rol 'admin'.")
    c.drawString(70, 550, "- Se modificó guestMiddleware.js para que los admins no sean redirigidos")
    c.drawString(70, 530, "  y puedan usar el formulario de registro.")
    c.drawString(70, 510, "- Se actualizó la vista register.ejs para mostrar un selector de rol a los admins.")
    c.drawString(70, 490, "- Se creó el seeder para 'Admin Thermo' (admin@thermo.com).")
    
    c.save()

def create_pptx(filepath):
    prs = Presentation()
    
    # Slide 1: Título
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "Corrección de Roles de Usuario"
    subtitle.text = "Explicación del problema y la solución implementada\nProyecto Pretty Thermo"
    
    # Slide 2: El problema
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    title_shape2 = shapes2.title
    body_shape2 = shapes2.placeholders[1]
    title_shape2.text = "¿Qué hiciste mal?"
    
    tf = body_shape2.text_frame
    p = tf.add_paragraph()
    p.text = "El formulario de registro y la sesión no manejaban correctamente el rol de usuario:"
    p.level = 0
    p2 = tf.add_paragraph()
    p2.text = "La sesión req.session.userLogged no guardaba la propiedad 'role'."
    p2.level = 1
    p3 = tf.add_paragraph()
    p3.text = "El proceso de registro delegaba el rol a los valores predeterminados sin validación estricta para prevenir la manipulación."
    p3.level = 1

    # Slide 3: La solución
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    title_shape3 = shapes3.title
    body_shape3 = shapes3.placeholders[1]
    title_shape3.text = "¿Cómo se corrigió?"
    
    tf3 = body_shape3.text_frame
    p4 = tf3.add_paragraph()
    p4.text = "Correcciones principales:"
    p4.level = 0
    p5 = tf3.add_paragraph()
    p5.text = "Al registrar usuarios públicos, el controlador fuerza estrictamente el rol a 'user'."
    p5.level = 1
    p6 = tf3.add_paragraph()
    p6.text = "El middleware de invitados ahora permite que los Administradores accedan a /register."
    p6.level = 1
    p7 = tf3.add_paragraph()
    p7.text = "En la vista, si el usuario logueado es Admin, se despliega un select para elegir 'Usuario' o 'Administrador'."
    p7.level = 1
    p8 = tf3.add_paragraph()
    p8.text = "Se creó un seeder para generar de forma segura al Admin Thermo por defecto."
    p8.level = 1
    
    prs.save(filepath)

if __name__ == "__main__":
    out_dir = r"c:\Users\alehx\OneDrive\Desktop\Digital House"
    if not os.path.exists(out_dir):
        os.makedirs(out_dir)
    
    pdf_path = os.path.join(out_dir, "reporte_cambios_roles.pdf")
    pptx_path = os.path.join(out_dir, "presentacion_cambios_roles.pptx")
    
    create_pdf(pdf_path)
    create_pptx(pptx_path)
    print("Reportes generados con exito en", out_dir)
